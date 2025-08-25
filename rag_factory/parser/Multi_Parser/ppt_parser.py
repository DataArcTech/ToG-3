import logging
from io import BytesIO
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import os


class PptParser:

    def __get_bulleted_text(self, paragraph):
        is_bulleted = bool(paragraph._p.xpath("./a:pPr/a:buChar")) or bool(paragraph._p.xpath("./a:pPr/a:buAutoNum")) or bool(paragraph._p.xpath("./a:pPr/a:buBlip"))
        if is_bulleted:
            return f"{'  '* paragraph.level}.{paragraph.text}"
        else:
            return paragraph.text
        

    def __extract(self, shape, output_dir):
        try:
            if hasattr(shape, 'has_text_frame') and shape.has_text_frame:
                text_frame = shape.text_frame
                texts = []
                for paragraph in text_frame.paragraphs:
                    if paragraph.text.strip():
                        texts.append(self.__get_bulleted_text(paragraph))
                return "text","".join(texts)
            try:
                shape_type = shape.shape_type
            except NotImplementedError:
                # If shape_type is not available, try to get text content
                if hasattr(shape, 'text'):
                    return 'text', shape.text.strip()
                return "",""

            # Handle table
            if shape_type == MSO_SHAPE_TYPE.TABLE:
                tb = shape.table
                html = "<table>\n"
                html += "  <tr>" + "".join(
                    f"<th>{tb.cell(0, j).text.strip()}</th>" for j in range(len(tb.columns))
                ) + "</tr>\n"
                for i in range(1, len(tb.rows)):
                    html += "  <tr>" + "".join(
                        f"<td>{tb.cell(i, j).text.strip()}</td>" for j in range(len(tb.columns))
                    ) + "</tr>\n"
                html += "</table>\n"
                return "table", html

            # Handle group shape
            if shape_type == MSO_SHAPE_TYPE.GROUP:
                texts = []
                for p in sorted(shape.shapes, key=lambda x: (x.top // 10, x.left)):
                    st,t = self.__extract(p)
                    if t:
                        texts.append(t.strip())
                return "group","\n".join(texts)
            # Handle image shape
            if shape_type == MSO_SHAPE_TYPE.PICTURE:
                import hashlib
                image = shape.image
                sha1 = hashlib.sha1(image.blob).hexdigest()
                ext = image.ext
                img_dir = os.path.join(output_dir, "image")
                if not os.path.exists(img_dir):
                    os.makedirs(img_dir)
                filename = os.path.join(img_dir, f"{sha1}.{ext}")
                with open(filename, "wb") as f:
                    f.write(image.blob)
                return "image", f"![image]({filename})"

            return "",""

        except Exception as e:
            logging.error(f"Error processing shape: {str(e)}")
            return "",""
        
    def _parse_pptx(self, fnm, from_page=0, to_page=None, output_dir=None):
        ppt = Presentation(fnm) if isinstance(fnm, str) else Presentation(BytesIO(fnm))
        file_name = os.path.basename(fnm) if isinstance(fnm, str) else "input.pptx"
        self.total_page = len(ppt.slides)
        if to_page is None:
            to_page = self.total_page

        md_slides = []
        for i, slide in enumerate(ppt.slides):
            if i < from_page:
                continue
            if i >= to_page:
                break
            slide_md = []
            for shape in sorted(
                slide.shapes,
                key=lambda x: ((x.top if x.top is not None else 0) // 10, x.left if x.left is not None else 0)
            ):
                try:
                    shape_type, txt = self.__extract(shape, output_dir)
                    if txt:
                        if shape_type == "text":
                            slide_md.append(txt + "\n")
                        elif shape_type == "table":
                            slide_md.append(txt + "\n") 
                        elif shape_type == "group":
                            slide_md.append(txt + "\n")
                        elif shape_type == "image":
                            slide_md.append(txt + "\n")
                except Exception as e:
                    logging.exception(e)

            slide_text = f"## Slide {i+1}\n" + "\n".join(slide_md)
            md_slides.append(slide_text)

        return md_slides
    
    def __call__(self, fnm, from_page=0, to_page=None, output_dir="./output"):
        if isinstance(fnm, str):
            file_name = os.path.basename(fnm).rsplit(".", 1)[0]
        elif isinstance(fnm, bytes):
            prefix_bytes = fnm[:10]
            prefix_hex = prefix_bytes.hex()
            # prefix_b64 = base64.urlsafe_b64encode(prefix_bytes).decode('utf-8')
            file_name= f"{prefix_hex}"
        else:
            file_name="input_pptx" 

        output_path = os.path.join(output_dir, file_name+"_markdown.md")
        md_slides = self._parse_pptx(fnm, from_page=from_page, to_page=to_page, output_dir=os.path.join(output_dir, file_name) )
        with open(output_path, 'w', encoding='utf-8') as output_f:
            for item in md_slides:
                output_f.write(item)
                # output_f.write("/n/n")
    
if __name__ == "__main__":
    parser = PptParser()
    parser('/home/yangcehao/doc_analysis/MultiTypeParser/12345.pptx',
           output_dir='/home/yangcehao/doc_analysis/MultiTypeParser/test')
